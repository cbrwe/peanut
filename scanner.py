"""
Peanut Scanner Engine
====================
Content-aware duplicate finder using perceptual hashing.
Streams results progressively via generator.

──────────────────────────────────────────────────────────────────────────
                    For whoever's reading this far.
──────────────────────────────────────────────────────────────────────────

The gap between what you're capable of and what you've been able to show
people is the most painful distance a person can live in.

It's not imposter syndrome. It's the opposite. You know exactly what's in
there. You've seen it in flashes. The problem isn't doubt. The problem is
that the gap is real and closing it requires conditions that nobody handed
you and you're still trying to build them yourself, from scratch, while
also just trying to get through the week.

──────────────────────────────────────────────────────────────────────────
"""

import hashlib
import json
import os
import re
import sqlite3
import subprocess
import tempfile
import threading
import time
from collections import defaultdict
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Generator, Optional

import imagehash
import numpy as np
from PIL import Image, ImageFilter, ExifTags
from tqdm import tqdm

# Register HEIC/HEIF support in PIL if pillow-heif is installed.
# This lets PIL natively decode iPhone photos (.heic) and high-efficiency
# files (.heif) for both phash + thumbnail. If pillow-heif isn't available,
# we silently fall through; make_thumbnail has an ffmpeg fallback that
# covers HEIC + RAW formats.
try:
    import pillow_heif
    pillow_heif.register_heif_opener()
except ImportError:
    pass

# ── Config ────────────────────────────────────────────────────────────────

IMAGE_EXTS = {
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".tif",
    ".webp", ".heic", ".heif", ".avif", ".ico",
    ".raw", ".cr2", ".nef", ".arw", ".dng", ".orf", ".rw2",  # RAW camera
    ".psd", ".svg",
    # Additional image formats (audit)
    ".jp2", ".j2k", ".jpf", ".jpx",   # JPEG 2000 family
    ".jxl",                            # JPEG XL
    ".jfif", ".pjpeg", ".pjp",        # JPEG variants
    ".dds",                            # DirectDraw Surface (game textures)
    ".tga",                            # Truevision TGA
    ".pbm", ".pgm", ".ppm", ".pnm",   # Netpbm
    ".cr3", ".raf", ".pef", ".srw",   # More RAW (Canon, Fuji, Pentax, Samsung)
}
VIDEO_EXTS = {
    ".mp4", ".mkv", ".avi", ".mov", ".wmv", ".flv", ".webm",
    ".m4v", ".mpg", ".mpeg", ".3gp", ".ts", ".mts", ".vob", ".ogv",
    # Additional video formats (audit)
    ".mxf",                            # Material Exchange Format (broadcast)
    ".f4v",                            # Flash Video (modern)
    ".divx",                           # DivX
    ".asf",                            # Advanced Systems Format
    ".rm", ".rmvb",                    # RealMedia
    ".m2ts",                           # Blu-ray transport stream
    ".m2v", ".mpv",                    # MPEG-2 video
    ".3g2",                            # 3GPP2
    ".mod", ".tod",                    # Camcorder formats
    ".dv",                             # Digital Video
}
AUDIO_EXTS = {
    ".mp3", ".wav", ".flac", ".aac", ".ogg", ".m4a", ".wma",
    ".alac", ".aiff", ".aif", ".opus", ".wv", ".ape",
    # Additional audio formats (audit)
    ".mid", ".midi",                   # MIDI
    ".amr",                            # Adaptive Multi-Rate (mobile recordings)
    ".ac3",                            # Dolby Digital
    ".dts",                            # DTS
    ".mka",                            # Matroska audio
    ".tta",                            # The True Audio (lossless)
    ".au", ".snd",                     # Sun/NeXT audio
    ".dsf", ".dff",                    # DSD (high-res audio)
    ".ra",                             # RealAudio
    ".caf",                            # Core Audio Format (Apple)
    ".voc",                            # Creative Voice
    ".m4b",                            # Audiobook (Apple)
    ".oga",                            # Ogg audio
}
LOSSLESS_EXTS = {".png", ".bmp", ".tiff", ".tif", ".gif", ".raw", ".cr2",
                 ".nef", ".arw", ".dng", ".orf", ".rw2", ".cr3", ".raf",
                 ".pef", ".srw", ".flac", ".wav", ".alac", ".aiff", ".aif",
                 ".wv", ".ape", ".tta", ".dsf", ".dff",
                 ".jxl", ".jp2", ".j2k", ".tga", ".pbm", ".pgm", ".ppm", ".pnm"}
LOSSY_EXTS = {".jpg", ".jpeg", ".jfif", ".pjpeg", ".pjp", ".webp", ".avif",
              ".heic", ".heif", ".mp3", ".aac", ".ogg", ".oga",
              ".m4a", ".m4b", ".wma", ".opus", ".amr", ".ac3", ".dts", ".ra"}

IMAGE_THRESHOLD = 10
VIDEO_THRESHOLD = 12
AUDIO_THRESHOLD = 8  # Audio fingerprint hamming distance
DOC_THRESHOLD = 6
VIDEO_KEYFRAMES = 5
MIN_FILE_SIZE = 1024

DOC_EXTS = {
    ".txt", ".md", ".rtf", ".log",
    ".csv", ".tsv",
    ".docx", ".doc",
    ".pptx", ".ppt",
    ".xlsx", ".xls",
    ".pdf",
    ".json", ".xml", ".html", ".htm",
    ".odt", ".ods", ".odp",  # LibreOffice
    ".epub",
    ".yaml", ".yml", ".toml", ".ini", ".cfg",
    ".sql", ".tex",
    # Additional document formats (audit)
    ".pages", ".numbers", ".key",      # Apple iWork
    ".mobi", ".azw", ".azw3",          # Kindle
    ".fb2", ".djvu",                   # Other ebook formats
    ".srt", ".vtt", ".sub", ".ass", ".ssa",  # Subtitles
    ".env", ".conf", ".config", ".properties",  # Config files
    ".bib",                            # BibTeX
    ".diff", ".patch",                 # Diffs
    ".org",                            # Org-mode
    ".markdown", ".mdown", ".mkd",     # Markdown variants
    ".rst",                            # reStructuredText
    ".adoc", ".asciidoc",              # AsciiDoc
}

CODE_EXTS = {
    ".py", ".js", ".ts", ".jsx", ".tsx", ".css", ".scss", ".sass",
    ".java", ".cpp", ".c", ".h", ".hpp", ".cs",
    ".swift", ".go", ".rs", ".rb", ".php",
    ".sh", ".bash", ".zsh", ".fish",
    ".r", ".m", ".lua", ".pl", ".kt",
    ".vue", ".svelte", ".astro",
    # Additional code formats (audit)
    ".dart",                           # Dart / Flutter
    ".scala", ".sc",                   # Scala
    ".clj", ".cljs", ".cljc", ".edn", # Clojure
    ".ex", ".exs",                     # Elixir
    ".hs", ".lhs",                     # Haskell
    ".zig",                            # Zig
    ".nim", ".nims",                   # Nim
    ".jl",                             # Julia
    ".f", ".f90", ".f95", ".f03",     # Fortran
    ".erl", ".hrl",                    # Erlang
    ".ml", ".mli",                     # OCaml / SML
    ".coffee",                         # CoffeeScript
    ".groovy", ".gradle",              # Groovy
    ".pas", ".pp",                     # Pascal
    ".d",                              # D
    ".cr",                             # Crystal
    ".v",                              # Verilog / V-lang
    ".rkt",                            # Racket
    ".lisp", ".cl", ".el",            # Lisp variants / Emacs Lisp
    ".asm", ".s",                      # Assembly
    ".ps1",                            # PowerShell
    ".vb", ".vbs",                     # VB
}

ARCHIVE_EXTS = {
    ".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".xz",
    ".tgz", ".tar.gz",
    # Additional archive formats (audit)
    ".tbz2", ".tar.bz2",               # tar+bzip2
    ".txz", ".tar.xz",                 # tar+xz
    ".tzst", ".tar.zst",               # tar+zstd
    ".lz", ".lzma", ".lzo",            # LZ family
    ".zst",                            # Zstandard
    ".cab",                            # Cabinet
    ".arj", ".lha", ".lzh",            # Older formats
    ".cpio",                           # cpio
    ".sit", ".sitx",                   # StuffIt
}

# OTHER_EXTS — files we recognize and exact-dedupe but don't deeply analyze.
# Includes fonts, disk images, executables, 3D/CAD, databases, serialization formats.
OTHER_EXTS = {
    # Fonts
    ".ttf", ".otf", ".woff", ".woff2", ".eot", ".ttc",
    # Disk images
    ".iso", ".dmg", ".img", ".vhd", ".vhdx", ".vmdk", ".qcow2",
    ".bin", ".cue", ".nrg", ".mdf", ".mds",
    # Executables / installers / packages
    ".exe", ".msi", ".deb", ".pkg", ".rpm", ".appimage", ".apk", ".ipa",
    ".dll", ".so", ".dylib",
    ".jar", ".war", ".ear",            # Java archives (executable)
    # 3D / CAD
    ".blend", ".stl", ".obj", ".fbx", ".dae", ".gltf", ".glb",
    ".dwg", ".dxf", ".skp", ".3ds", ".max", ".c4d",
    ".step", ".stp", ".iges", ".igs", ".x3d",
    # Databases
    ".db", ".sqlite", ".sqlite3", ".db3",
    ".mdb", ".accdb",                  # MS Access
    # Serialization / data formats
    ".pkl", ".pickle",                 # Python pickle
    ".parquet", ".feather", ".arrow",  # Columnar data
    ".avro",                           # Apache Avro
    ".pb", ".protobuf",                # Protocol Buffers
    ".npy", ".npz",                    # NumPy
    ".h5", ".hdf5", ".nc",             # HDF5 / NetCDF
    ".mat",                            # MATLAB
    ".rds", ".rda",                    # R data
    ".sav",                            # SPSS
    ".dta",                            # Stata
}

# PLAIN_TEXT_DOC_EXTS — formats that should be read as raw text for SimHash.
# These were silently failing text extraction before the audit (their bytes
# were skipped, producing empty text_hash and no similarity matches).
PLAIN_TEXT_DOC_EXTS = {
    # Already-supported plain text (kept for clarity / single source of truth)
    ".txt", ".md", ".log", ".rtf",
    ".csv", ".tsv",
    ".json", ".xml", ".html", ".htm",
    # Newly-supported plain text formats
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".config",
    ".env", ".properties",
    ".sql", ".tex", ".bib",
    ".srt", ".vtt", ".sub", ".ass", ".ssa",  # Subtitles
    ".diff", ".patch",
    ".org",
    ".markdown", ".mdown", ".mkd",
    ".rst",
    ".adoc", ".asciidoc",
}

# File type category map for UI filtering
FILE_CATEGORIES = {
    "images": IMAGE_EXTS,
    "videos": VIDEO_EXTS,
    "audio": AUDIO_EXTS,
    "documents": DOC_EXTS,
    "code": CODE_EXTS,
    "archives": ARCHIVE_EXTS,
    "other": OTHER_EXTS,
}

ALL_EXTS = (IMAGE_EXTS | VIDEO_EXTS | AUDIO_EXTS | DOC_EXTS
            | CODE_EXTS | ARCHIVE_EXTS | OTHER_EXTS)

# ── Source Detection Patterns ─────────────────────────────────────────────

SOURCE_PATH_PATTERNS = [
    (r"whatsapp", "WhatsApp"),
    (r"telegram", "Telegram"),
    (r"signal", "Signal"),
    (r"imessage|messages.*attachments", "iMessage"),
    (r"screenshots?", "Screenshot"),
    (r"screen.?record", "Screen Recording"),
    (r"downloads?", "Download"),
    (r"desktop", "Desktop Save"),
    (r"google.?photos?|gphotos", "Google Photos"),
    (r"icloud|photostream", "iCloud"),
    (r"dropbox", "Dropbox"),
    (r"onedrive", "OneDrive"),
    (r"slack", "Slack"),
    (r"instagram", "Instagram"),
    (r"facebook|fb", "Facebook"),
    (r"twitter|x\.com", "Twitter/X"),
    (r"reddit", "Reddit"),
    (r"lightroom|lrcat", "Lightroom"),
    (r"photoshop|\.psd", "Photoshop"),
    (r"snapseed", "Snapseed"),
    (r"vsco", "VSCO"),
    (r"airdrop", "AirDrop"),
    (r"dcim", "Camera Roll"),
]

SOURCE_FILENAME_PATTERNS = [
    (r"^IMG_\d{4}", "iPhone Camera"),
    (r"^DSC_?\d{4}", "Nikon DSLR"),
    (r"^_DSC\d{4}", "Sony Camera"),
    (r"^DSCF\d{4}", "Fujifilm"),
    (r"^P\d{7}", "Panasonic"),
    (r"^GOPR\d{4}|^GX\d{6}", "GoPro"),
    (r"^DJI_\d{4}", "DJI Drone"),
    (r"^Screenshot", "Screenshot"),
    (r"^Screen Recording", "Screen Recording"),
    (r"^Photo \d", "macOS Photos Export"),
    (r"^IMG-\d{8}-WA", "WhatsApp"),
    (r"^signal-\d", "Signal"),
    (r"^PXL_\d{8}", "Google Pixel"),
    (r"^Samsung_\d", "Samsung Camera"),
]


# ── Vectorization helpers ────────────────────────────────────────────────

# Popcount lookup table for uint8 (256 entries) — used for vectorized
# Hamming-distance computation on packed bit arrays.
_POPCOUNT_LUT = np.array([bin(i).count("1") for i in range(256)], dtype=np.uint8)


def _hex_to_bits_array(hex_strings: list, expected_hex_len: int = None) -> np.ndarray:
    """
    Convert a list of hex strings to a packed-byte numpy array suitable for
    vectorized Hamming-distance computation.

    Returns a (n, n_bytes) uint8 array where each row is the byte-packed
    representation of one hash. Hex strings of unexpected length are
    zero-padded or truncated to ``expected_hex_len`` so all rows align.

    If ``expected_hex_len`` is None, the modal length of the inputs is used.
    """
    if not hex_strings:
        return np.zeros((0, 0), dtype=np.uint8)

    if expected_hex_len is None:
        # Use the most common length
        from collections import Counter
        expected_hex_len = Counter(len(s) for s in hex_strings).most_common(1)[0][0]
    # Force to even (each byte = 2 hex chars)
    if expected_hex_len % 2:
        expected_hex_len += 1
    n_bytes = expected_hex_len // 2

    out = np.zeros((len(hex_strings), n_bytes), dtype=np.uint8)
    for i, s in enumerate(hex_strings):
        # Normalize: pad with zeros or truncate
        if len(s) < expected_hex_len:
            s = s.zfill(expected_hex_len)
        elif len(s) > expected_hex_len:
            s = s[:expected_hex_len]
        try:
            out[i] = np.frombuffer(bytes.fromhex(s), dtype=np.uint8)
        except ValueError:
            # Bad hex — leave row as zeros (will be high-distance from anything)
            pass
    return out


def _vectorized_pairs(packed: np.ndarray, threshold: float,
                      divisor: int = 1) -> list:
    """
    Find all index pairs (i, j) with i < j whose Hamming distance is <= threshold.

    ``packed`` is the (n, n_bytes) uint8 array from ``_hex_to_bits_array``,
    optionally with multiple segments concatenated horizontally so the total
    distance is the sum across segments. ``divisor`` is the number of segments
    used to convert summed distance into a per-segment average.

    Returns a list of (i, j, distance) tuples where distance is the float
    average per segment (summed bits / divisor).
    """
    n = packed.shape[0]
    if n < 2:
        return []

    pairs = []
    # For each row i, vectorize XOR + popcount against all rows j > i
    for i in range(n - 1):
        xor = packed[i + 1:] ^ packed[i]                       # (n-i-1, n_bytes) uint8
        # _POPCOUNT_LUT lookup gives bit counts per byte; sum across bytes.
        # Use int64 to prevent overflow for long hashes (e.g. videos with 5x256-bit segments).
        dist = _POPCOUNT_LUT[xor].sum(axis=1, dtype=np.int64)   # (n-i-1,)
        if divisor > 1:
            # Average distance per segment
            below = np.nonzero(dist <= threshold * divisor)[0]
            for k in below:
                pairs.append((i, i + 1 + int(k), float(dist[k]) / divisor))
        else:
            below = np.nonzero(dist <= threshold)[0]
            for k in below:
                pairs.append((i, i + 1 + int(k), int(dist[k])))
    return pairs


# ── Data Classes ──────────────────────────────────────────────────────────

@dataclass
class FileInfo:
    path: str
    size: int
    file_type: str
    md5: str = ""
    phash: str = ""
    dimensions: str = ""
    width: int = 0
    height: int = 0
    duration: float = 0.0
    modified: float = 0.0
    thumbnail: str = ""
    # Quality
    sharpness: float = 0.0
    bits_per_pixel: float = 0.0
    quality_score: int = 0  # 0-100
    compression: str = ""  # "lossless" | "lossy" | "raw" | ""
    # Metadata
    exif: dict = field(default_factory=dict)
    source: str = ""  # Detected origin
    # Text content (documents)
    text_hash: str = ""  # SimHash
    word_count: int = 0
    text_preview: str = ""  # First ~300 chars for UI preview
    # Comparison
    similarity: float = 100.0

    def to_dict(self):
        d = asdict(self)
        d["name"] = Path(self.path).name
        d["dir"] = str(Path(self.path).parent)
        d["ext"] = Path(self.path).suffix.lower()
        d["megapixels"] = round(self.width * self.height / 1_000_000, 1) if self.width else 0
        return d


@dataclass
class DuplicateGroup:
    group_id: int
    match_type: str
    files: list = field(default_factory=list)
    recommended_keep: str = ""

    def to_dict(self):
        return {
            "group_id": self.group_id,
            "match_type": self.match_type,
            "files": [f.to_dict() for f in self.files],
            "recommended_keep": self.recommended_keep,
            "total_size": sum(f.size for f in self.files),
            "recoverable_size": sum(
                f.size for f in self.files if f.path != self.recommended_keep
            ),
            "file_count": len(self.files),
        }


# ── Utility Functions ─────────────────────────────────────────────────────

def md5_hash(path: str) -> str:
    h = hashlib.md5()
    try:
        with open(path, "rb") as f:
            while chunk := f.read(8192):
                h.update(chunk)
        return h.hexdigest()
    except (OSError, PermissionError):
        return ""


def image_phash(path: str) -> Optional[str]:
    try:
        img = Image.open(path).convert("RGB")
        return str(imagehash.phash(img, hash_size=16))
    except Exception:
        return None


def image_dimensions(path: str) -> tuple:
    try:
        img = Image.open(path)
        return img.width, img.height
    except Exception:
        return 0, 0


def compute_sharpness(path: str) -> float:
    """Sharpness via variance of edge-detected image."""
    try:
        img = Image.open(path).convert("L")
        # Resize for consistent scoring and speed
        img.thumbnail((512, 512), Image.LANCZOS)
        edges = img.filter(ImageFilter.FIND_EDGES)
        arr = np.array(edges, dtype=np.float64)
        return float(np.var(arr))
    except Exception:
        return 0.0


def compute_quality_score(f: "FileInfo") -> int:
    """Composite quality score 0-100."""
    score = 0.0
    # Resolution component (max 40 points)
    mp = (f.width * f.height) / 1_000_000 if f.width else 0
    score += min(40, mp * 4)  # 10MP = 40pts
    # Sharpness component (max 30 points)
    score += min(30, f.sharpness / 100)
    # Bits per pixel (max 20 points) - higher = less compressed
    score += min(20, f.bits_per_pixel * 2.5)
    # Lossless bonus (10 points)
    if f.compression == "lossless":
        score += 10
    return min(100, max(0, int(score)))


def extract_exif(path: str) -> dict:
    """Extract key EXIF fields from an image."""
    result = {}
    try:
        img = Image.open(path)
        exif_data = img._getexif()
        if not exif_data:
            return result

        tag_names = {v: k for k, v in ExifTags.TAGS.items()}
        wanted = {
            "Make": "camera_make",
            "Model": "camera_model",
            "DateTime": "date_taken",
            "DateTimeOriginal": "date_taken",
            "Software": "software",
            "ImageWidth": "exif_width",
            "ImageLength": "exif_height",
            "ExifImageWidth": "exif_width",
            "ExifImageHeight": "exif_height",
            "ISOSpeedRatings": "iso",
            "FNumber": "aperture",
            "ExposureTime": "shutter_speed",
            "FocalLength": "focal_length",
            "LensModel": "lens",
            "GPSInfo": "gps",
        }

        for tag_id, value in exif_data.items():
            tag_name = ExifTags.TAGS.get(tag_id, "")
            if tag_name in wanted:
                key = wanted[tag_name]
                if key == "gps":
                    result["has_gps"] = True
                    continue
                if key == "aperture" and hasattr(value, "numerator"):
                    value = round(float(value), 1)
                elif key == "shutter_speed" and hasattr(value, "numerator"):
                    if value.numerator and value.denominator:
                        if float(value) < 1:
                            value = f"1/{int(value.denominator/value.numerator)}"
                        else:
                            value = f"{float(value):.1f}"
                elif key == "focal_length" and hasattr(value, "numerator"):
                    value = f"{round(float(value))}mm"
                if key not in result:  # Don't overwrite DateTimeOriginal with DateTime
                    result[key] = str(value) if not isinstance(value, (int, float)) else value
    except Exception:
        pass
    return result


def detect_source(path: str, exif: dict) -> str:
    """Detect the likely origin of a file."""
    path_lower = path.lower()
    name = Path(path).name

    # Check EXIF software first
    sw = exif.get("software", "").lower()
    if sw:
        for pattern, source in [
            ("instagram", "Instagram"), ("snapchat", "Snapchat"),
            ("vsco", "VSCO"), ("lightroom", "Lightroom"),
            ("photoshop", "Photoshop"), ("snapseed", "Snapseed"),
        ]:
            if pattern in sw:
                return source

    # Check filename patterns
    for pattern, source in SOURCE_FILENAME_PATTERNS:
        if re.match(pattern, name, re.IGNORECASE):
            return source

    # Check path patterns
    for pattern, source in SOURCE_PATH_PATTERNS:
        if re.search(pattern, path_lower):
            return source

    # Camera make from EXIF
    make = exif.get("camera_make", "")
    model = exif.get("camera_model", "")
    if make or model:
        return f"{make} {model}".strip()

    return ""


def detect_compression(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext in LOSSLESS_EXTS:
        return "lossless"
    elif ext in LOSSY_EXTS:
        return "lossy"
    elif ext in {".heic", ".heif"}:
        return "lossy"  # HEIC is typically lossy
    elif ext in {".tiff", ".tif"}:
        return "lossless"
    return ""


# ── Text Extraction ───────────────────────────────────────────────────────

def extract_text_from_file(path: str) -> str:
    """Extract text content from any supported document type. Times out after 10s."""
    import threading

    result = [""]
    def _extract():
        try:
            result[0] = _extract_text_inner(path)
        except Exception:
            result[0] = ""

    t = threading.Thread(target=_extract, daemon=True)
    t.start()
    t.join(timeout=10)
    return result[0]


def _extract_text_inner(path: str) -> str:
    """Inner text extraction (called with timeout wrapper)."""
    ext = Path(path).suffix.lower()
    try:
        # Plain-text formats (txt, md, csv, json, xml, html, yaml, toml, ini,
        # sql, tex, srt, vtt, env, etc.). Read as bytes, decode with fallbacks.
        if ext in PLAIN_TEXT_DOC_EXTS:
            for enc in ("utf-8", "latin-1", "cp1252"):
                try:
                    with open(path, "r", encoding=enc, errors="replace") as f:
                        return f.read(500_000)
                except (UnicodeDecodeError, OSError):
                    continue
            return ""
        elif ext == ".docx":
            try:
                import docx
                doc = docx.Document(path)
                parts = [p.text for p in doc.paragraphs if p.text.strip()]
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                parts.append(cell.text)
                return "\n".join(parts)
            except ImportError:
                import zipfile, re as _re
                with zipfile.ZipFile(path, "r") as zf:
                    data = zf.read("word/document.xml").decode("utf-8", errors="replace")
                    return _re.sub(r"\s+", " ", _re.sub(r"<[^>]+>", " ", data)).strip()
        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(path)
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            parts.append(shape.text)
                return "\n".join(parts)
            except ImportError:
                import zipfile, re as _re
                with zipfile.ZipFile(path, "r") as zf:
                    parts = []
                    for name in zf.namelist():
                        if name.startswith("ppt/slides/") and name.endswith(".xml"):
                            data = zf.read(name).decode("utf-8", errors="replace")
                            text = _re.sub(r"\s+", " ", _re.sub(r"<[^>]+>", " ", data)).strip()
                            if text:
                                parts.append(text)
                    return "\n".join(parts)
        elif ext in {".xlsx", ".xls"}:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
                parts = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(max_row=5000, values_only=True):
                        cells = [str(c) for c in row if c is not None]
                        if cells:
                            parts.append(" ".join(cells))
                wb.close()
                return "\n".join(parts)
            except ImportError:
                import zipfile, re as _re
                with zipfile.ZipFile(path, "r") as zf:
                    if "xl/sharedStrings.xml" in zf.namelist():
                        data = zf.read("xl/sharedStrings.xml").decode("utf-8", errors="replace")
                        return _re.sub(r"\s+", " ", _re.sub(r"<[^>]+>", " ", data)).strip()
        elif ext == ".pdf":
            try:
                import PyPDF2
                with open(path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    parts = []
                    for page in reader.pages[:200]:
                        text = page.extract_text()
                        if text and text.strip():
                            parts.append(text)
                    return "\n".join(parts)
            except ImportError:
                pass
        elif ext in {".odt", ".ods", ".odp"}:
            # OpenDocument formats are zip archives containing content.xml.
            import zipfile, re as _re
            try:
                with zipfile.ZipFile(path, "r") as zf:
                    if "content.xml" in zf.namelist():
                        data = zf.read("content.xml").decode("utf-8", errors="replace")
                        return _re.sub(r"\s+", " ", _re.sub(r"<[^>]+>", " ", data)).strip()
            except (zipfile.BadZipFile, OSError):
                pass
        elif ext == ".epub":
            # EPUB is a zip with HTML/XHTML content files.
            import zipfile, re as _re
            try:
                with zipfile.ZipFile(path, "r") as zf:
                    parts = []
                    # Read HTML/XHTML files in package order if possible,
                    # else iterate alphabetically (good enough for SimHash).
                    text_files = sorted(
                        n for n in zf.namelist()
                        if n.lower().endswith((".html", ".xhtml", ".htm"))
                    )
                    total_len = 0
                    for name in text_files:
                        try:
                            data = zf.read(name).decode("utf-8", errors="replace")
                        except (KeyError, OSError):
                            continue
                        text = _re.sub(r"\s+", " ", _re.sub(r"<[^>]+>", " ", data)).strip()
                        if text:
                            parts.append(text)
                            total_len += len(text)
                            # Cap total extracted text to keep timeout-bound
                            if total_len > 500_000:
                                break
                    return "\n".join(parts)
            except (zipfile.BadZipFile, OSError):
                pass
    except Exception:
        pass
    return ""


def compute_simhash(text: str, hash_bits: int = 128) -> str:
    """Locality-sensitive hash for text. Similar docs produce similar hashes."""
    if not text or len(text.strip()) < 20:
        return ""
    words = text.lower().split()
    if len(words) < 3:
        return ""
    shingles = [" ".join(words[i:i+3]) for i in range(len(words) - 2)]
    if not shingles:
        return ""
    v = [0] * hash_bits
    for shingle in shingles:
        h = int(hashlib.md5(shingle.encode("utf-8", errors="replace")).hexdigest(), 16)
        for i in range(hash_bits):
            if h & (1 << i):
                v[i] += 1
            else:
                v[i] -= 1
    fp = 0
    for i in range(hash_bits):
        if v[i] > 0:
            fp |= (1 << i)
    return format(fp, f"0{hash_bits // 4}x")


def simhash_distance(h1: str, h2: str) -> int:
    """Hamming distance between two SimHash hex strings."""
    try:
        return bin(int(h1, 16) ^ int(h2, 16)).count("1")
    except (ValueError, TypeError):
        return 999


def ffmpeg_ok() -> bool:
    try:
        subprocess.run(["ffmpeg", "-version"], capture_output=True, timeout=5)
        return True
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return False


def video_duration(path: str) -> float:
    try:
        r = subprocess.run(
            ["ffprobe", "-v", "quiet", "-print_format", "json", "-show_format", path],
            capture_output=True, text=True, timeout=30,
        )
        return float(json.loads(r.stdout).get("format", {}).get("duration", 0))
    except Exception:
        return 0.0


def video_dimensions(path: str) -> tuple:
    try:
        r = subprocess.run(
            ["ffprobe", "-v", "quiet", "-print_format", "json",
             "-show_streams", "-select_streams", "v:0", path],
            capture_output=True, text=True, timeout=30,
        )
        streams = json.loads(r.stdout).get("streams", [])
        if streams:
            return streams[0].get("width", 0), streams[0].get("height", 0)
    except Exception:
        pass
    return 0, 0


def video_phash(path: str) -> Optional[str]:
    """Composite perceptual hash from video keyframes."""
    try:
        dur = video_duration(path)
        if dur <= 0:
            return None
        hashes = []
        with tempfile.TemporaryDirectory() as tmp:
            timestamps = [dur * (i + 1) / (VIDEO_KEYFRAMES + 1) for i in range(VIDEO_KEYFRAMES)]
            for i, ts in enumerate(timestamps):
                out = os.path.join(tmp, f"f{i}.jpg")
                subprocess.run(
                    ["ffmpeg", "-ss", str(ts), "-i", path, "-vframes", "1", "-q:v", "2", "-y", out],
                    capture_output=True, timeout=30,
                )
                if os.path.exists(out):
                    img = Image.open(out).convert("RGB")
                    hashes.append(str(imagehash.phash(img, hash_size=16)))
        return ";".join(hashes) if hashes else None
    except Exception:
        return None


def hamming(h1: str, h2: str) -> int:
    try:
        return imagehash.hex_to_hash(h1) - imagehash.hex_to_hash(h2)
    except Exception:
        return 999


def video_hamming(h1: str, h2: str) -> float:
    f1, f2 = h1.split(";"), h2.split(";")
    dists = [hamming(f1[i], f2[i]) for i in range(min(len(f1), len(f2)))]
    return sum(dists) / len(dists) if dists else 999.0


def audio_duration(path: str) -> float:
    """Get audio duration via ffprobe."""
    try:
        r = subprocess.run(
            ["ffprobe", "-v", "quiet", "-print_format", "json", "-show_format", path],
            capture_output=True, text=True, timeout=15,
        )
        return float(json.loads(r.stdout).get("format", {}).get("duration", 0))
    except Exception:
        return 0.0


def audio_fingerprint(path: str) -> Optional[str]:
    """
    Generate audio fingerprint by converting to raw PCM samples,
    chunking into segments, and hashing each segment's energy distribution.
    Returns semicolon-separated hex hashes.
    """
    try:
        with tempfile.TemporaryDirectory() as tmp:
            raw_path = os.path.join(tmp, "audio.raw")
            # Convert to mono 8kHz raw PCM via ffmpeg
            subprocess.run(
                ["ffmpeg", "-i", path, "-ac", "1", "-ar", "8000",
                 "-f", "s16le", "-t", "30",  # first 30 seconds
                 "-y", raw_path],
                capture_output=True, timeout=20,
            )
            if not os.path.exists(raw_path) or os.path.getsize(raw_path) < 1600:
                return None

            data = np.fromfile(raw_path, dtype=np.int16)
            # Split into 10 segments and hash each
            seg_len = len(data) // 10
            if seg_len < 100:
                return None

            hashes = []
            for i in range(10):
                seg = data[i * seg_len:(i + 1) * seg_len].astype(np.float64)
                # Compute simple spectral energy signature
                fft = np.abs(np.fft.rfft(seg))
                # Bucket into 32 frequency bands
                bands = np.array_split(fft, 32)
                energies = [np.mean(b) for b in bands]
                # Binary hash: each band above/below median
                median = np.median(energies)
                bits = "".join("1" if e > median else "0" for e in energies)
                h = format(int(bits, 2), "08x")
                hashes.append(h)
            return ";".join(hashes)
    except Exception:
        return None


def audio_hash_distance(h1: str, h2: str) -> float:
    """Average hamming distance between audio fingerprint segments."""
    s1, s2 = h1.split(";"), h2.split(";")
    dists = []
    for i in range(min(len(s1), len(s2))):
        try:
            dists.append(bin(int(s1[i], 16) ^ int(s2[i], 16)).count("1"))
        except ValueError:
            dists.append(32)
    return sum(dists) / len(dists) if dists else 999.0


def archive_content_hash(path: str) -> str:
    """Hash an archive's file listing (names + sizes) for duplicate detection.

    Supports zip and tar variants (.tar, .tar.gz/.tgz, .tar.bz2/.tbz2,
    .tar.xz/.txz). For other archive formats (.rar, .7z, .zst, etc.) we
    fall back to a head+tail+size byte hash, which is stable but won't
    detect content-equivalent archives with different compression.
    """
    import zipfile
    import tarfile
    pl = path.lower()
    try:
        if pl.endswith(".zip"):
            with zipfile.ZipFile(path, "r") as zf:
                listing = sorted(f"{i.filename}:{i.file_size}" for i in zf.infolist()
                                 if not i.is_dir())
                return hashlib.md5("\n".join(listing).encode()).hexdigest()
        # tar and compressed-tar variants
        tar_suffixes = (".tar", ".tar.gz", ".tgz", ".tar.bz2", ".tbz2",
                        ".tar.xz", ".txz")
        if pl.endswith(tar_suffixes):
            # tarfile auto-detects gzip/bzip2/xz from magic bytes when mode='r'
            with tarfile.open(path, mode="r") as tf:
                listing = sorted(
                    f"{m.name}:{m.size}" for m in tf.getmembers() if m.isfile()
                )
                return hashlib.md5("\n".join(listing).encode()).hexdigest()
    except (zipfile.BadZipFile, tarfile.TarError, OSError, EOFError):
        pass
    except Exception:
        pass
    # For non-zip/non-tar archives (rar, 7z, zst, cab, etc.), or any archive
    # we couldn't open, hash size + first/last 4KB.
    try:
        with open(path, "rb") as f:
            head = f.read(4096)
            try:
                f.seek(-4096, 2)
                tail = f.read(4096)
            except OSError:
                tail = b""
            return hashlib.md5(head + tail).hexdigest()
    except Exception:
        return ""


def similarity_pct(dist: float, max_dist: float = 256.0) -> float:
    """Convert hamming distance to similarity percentage."""
    return round(max(0, (1 - dist / max_dist)) * 100, 1)


def make_thumbnail(path: str, thumb_dir: str, file_type: str) -> str:
    fhash = hashlib.md5(path.encode()).hexdigest()[:12]
    thumb_name = f"{fhash}.jpg"
    thumb_path = os.path.join(thumb_dir, thumb_name)
    if os.path.exists(thumb_path):
        return f"/thumb/{thumb_name}"
    try:
        if file_type == "image":
            # Try PIL first (fast, no subprocess)
            try:
                img = Image.open(path).convert("RGB")
                img.thumbnail((400, 400), Image.LANCZOS)
                img.save(thumb_path, "JPEG", quality=82)
            except Exception:
                # PIL can't decode HEIC, some RAW formats, or corrupt files.
                # Fall back to ffmpeg, which handles HEIC + most RAW natively.
                try:
                    subprocess.run(
                        ["ffmpeg", "-i", path, "-vframes", "1",
                         "-vf", "scale=400:-1", "-q:v", "3", "-y", thumb_path],
                        capture_output=True, timeout=15,
                    )
                except (FileNotFoundError, subprocess.TimeoutExpired):
                    pass
        elif file_type == "video":
            dur = video_duration(path)
            seek = max(0, dur * 0.1) if dur > 0 else 1
            subprocess.run(
                ["ffmpeg", "-ss", str(seek), "-i", path, "-vframes", "1",
                 "-vf", "scale=400:-1", "-q:v", "3", "-y", thumb_path],
                capture_output=True, timeout=15,
            )
        return f"/thumb/{thumb_name}" if os.path.exists(thumb_path) else ""
    except Exception:
        return ""


def pick_best(files: list) -> str:
    """Pick file to keep: highest quality score, then resolution, then size."""
    if not files:
        return ""
    def score(f):
        # Prefer organized folders over Downloads/Desktop
        folder_bonus = 0
        pl = f.path.lower()
        if any(x in pl for x in ["/downloads", "/desktop", "/tmp", "/temp"]):
            folder_bonus = -10
        if any(x in pl for x in ["/photos", "/pictures", "/library", "/organized"]):
            folder_bonus = 10
        return (f.quality_score + folder_bonus, f.width * f.height, f.size, f.modified)
    return max(files, key=score).path


# ── Cache ─────────────────────────────────────────────────────────────────

class HashCache:
    def __init__(self, db_path: str):
        # check_same_thread=False — we share this connection across
        # the main thread, parallel workers, and Flask request threads.
        # timeout=30 — when SQLite hits a locked database, wait up to
        # 30 seconds for the lock to clear instead of failing instantly.
        # Real-world scans of 70k+ files can have brief periods where
        # the cache.db is being written by one thread while another
        # tries to read or write; without a timeout, that's an instant
        # "database is locked" error and the scan dies.
        self.conn = sqlite3.connect(db_path, check_same_thread=False, timeout=30.0)
        # Performance pragmas. WAL mode lets readers and writers coexist
        # without blocking each other. NORMAL synchronous trades a tiny
        # crash-safety reduction for big throughput gains (we can always
        # rebuild the cache by re-scanning if it gets corrupted). The
        # busy_timeout is belt-and-suspenders alongside the connect timeout.
        try:
            self.conn.execute("PRAGMA journal_mode=WAL")
            self.conn.execute("PRAGMA synchronous=NORMAL")
            self.conn.execute("PRAGMA busy_timeout=30000")
            self.conn.execute("PRAGMA temp_store=MEMORY")
        except Exception:
            # PRAGMAs failing is non-fatal — we just lose the perf boost.
            pass
        self._pending_writes = 0
        self._write_lock = threading.Lock()  # serialize writes from main + worker threads
        self.conn.execute("""
            CREATE TABLE IF NOT EXISTS file_hashes (
                path TEXT PRIMARY KEY,
                size INTEGER,
                mtime REAL,
                md5 TEXT,
                phash TEXT,
                sharpness REAL,
                exif_json TEXT,
                source TEXT,
                compression TEXT,
                width INTEGER,
                height INTEGER,
                text_hash TEXT DEFAULT '',
                word_count INTEGER DEFAULT 0,
                text_preview TEXT DEFAULT ''
            )
        """)
        # Lock table — tracks files the user has marked as protected from
        # deletion. Locks persist across scans (separate from file_hashes
        # so they survive cache invalidation when a file's mtime/size
        # changes).
        self.conn.execute("""
            CREATE TABLE IF NOT EXISTS locked_files (
                path TEXT PRIMARY KEY,
                locked_at TEXT
            )
        """)
        # Migrate old tables missing new columns
        try:
            self.conn.execute("ALTER TABLE file_hashes ADD COLUMN text_hash TEXT DEFAULT ''")
        except Exception:
            pass
        try:
            self.conn.execute("ALTER TABLE file_hashes ADD COLUMN word_count INTEGER DEFAULT 0")
        except Exception:
            pass
        try:
            self.conn.execute("ALTER TABLE file_hashes ADD COLUMN text_preview TEXT DEFAULT ''")
        except Exception:
            pass
        self.conn.commit()

    # ── Locks API ────────────────────────────────────────────────────────
    # These are deliberately separate from file_hashes so locks persist
    # even if a file's content/cache row gets invalidated.

    def lock(self, path: str) -> None:
        from datetime import datetime
        with self._write_lock:
            self.conn.execute(
                "INSERT OR REPLACE INTO locked_files (path, locked_at) VALUES (?, ?)",
                (path, datetime.utcnow().isoformat() + "Z"),
            )
            self.conn.commit()

    def unlock(self, path: str) -> None:
        with self._write_lock:
            self.conn.execute("DELETE FROM locked_files WHERE path=?", (path,))
            self.conn.commit()

    def is_locked(self, path: str) -> bool:
        row = self.conn.execute(
            "SELECT 1 FROM locked_files WHERE path=?", (path,)
        ).fetchone()
        return row is not None

    def list_locked(self) -> set:
        rows = self.conn.execute("SELECT path FROM locked_files").fetchall()
        return {r[0] for r in rows}

    def get(self, path: str, size: int, mtime: float) -> Optional[dict]:
        row = self.conn.execute(
            "SELECT md5, phash, sharpness, exif_json, source, compression, width, height, text_hash, word_count, text_preview FROM file_hashes WHERE path=? AND size=? AND mtime=?",
            (path, size, mtime),
        ).fetchone()
        if row:
            return {
                "md5": row[0], "phash": row[1], "sharpness": row[2],
                "exif": json.loads(row[3]) if row[3] else {},
                "source": row[4], "compression": row[5],
                "width": row[6], "height": row[7],
                "text_hash": row[8] or "", "word_count": row[9] or 0,
                "text_preview": row[10] or "",
            }
        return None

    def put(self, path: str, size: int, mtime: float, data: dict):
        """
        Insert or update a cache row, MERGING with any existing row instead of
        clobbering. This prevents image-only updates from wiping text fields
        (text_hash, word_count, text_preview) and vice versa.

        Commits are batched: every 100 puts we commit, and flush() can be
        called explicitly to commit any pending writes.
        """
        # Read existing row to preserve fields not present in `data`.
        # The whole read-modify-write is inside the write lock so two
        # concurrent put()s for the same path can't clobber each other.
        with self._write_lock:
            existing = self.conn.execute(
                "SELECT md5, phash, sharpness, exif_json, source, compression, "
                "width, height, text_hash, word_count, text_preview "
                "FROM file_hashes WHERE path=?",
                (path,),
            ).fetchone()
            if existing:
                cur = {
                    "md5": existing[0] or "",
                    "phash": existing[1] or "",
                    "sharpness": existing[2] or 0,
                    "exif_json": existing[3] or "{}",
                    "source": existing[4] or "",
                    "compression": existing[5] or "",
                    "width": existing[6] or 0,
                    "height": existing[7] or 0,
                    "text_hash": existing[8] or "",
                    "word_count": existing[9] or 0,
                    "text_preview": existing[10] or "",
                }
            else:
                cur = {
                    "md5": "", "phash": "", "sharpness": 0, "exif_json": "{}",
                    "source": "", "compression": "", "width": 0, "height": 0,
                    "text_hash": "", "word_count": 0, "text_preview": "",
                }
            # Apply only the fields present in `data` (with sensible "is set" tests)
            if "md5" in data and data["md5"]:                      cur["md5"] = data["md5"]
            if "phash" in data and data["phash"]:                  cur["phash"] = data["phash"]
            if "sharpness" in data and data["sharpness"]:          cur["sharpness"] = data["sharpness"]
            if "exif" in data and data["exif"]:                    cur["exif_json"] = json.dumps(data["exif"])
            if "source" in data and data["source"]:                cur["source"] = data["source"]
            if "compression" in data and data["compression"]:      cur["compression"] = data["compression"]
            if "width" in data and data["width"]:                  cur["width"] = data["width"]
            if "height" in data and data["height"]:                cur["height"] = data["height"]
            if "text_hash" in data and data["text_hash"]:          cur["text_hash"] = data["text_hash"]
            if "word_count" in data and data["word_count"]:        cur["word_count"] = data["word_count"]
            if "text_preview" in data and data["text_preview"]:    cur["text_preview"] = data["text_preview"]

            self.conn.execute(
                """INSERT OR REPLACE INTO file_hashes
                   (path, size, mtime, md5, phash, sharpness, exif_json, source, compression, width, height, text_hash, word_count, text_preview)
                   VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
                (path, size, mtime, cur["md5"], cur["phash"],
                 cur["sharpness"], cur["exif_json"],
                 cur["source"], cur["compression"],
                 cur["width"], cur["height"],
                 cur["text_hash"], cur["word_count"], cur["text_preview"]),
            )
            self._pending_writes += 1
            if self._pending_writes >= 100:
                self.conn.commit()
                self._pending_writes = 0

    def flush(self):
        """Commit any pending writes to disk."""
        with self._write_lock:
            if self._pending_writes > 0:
                self.conn.commit()
                self._pending_writes = 0

    def close(self):
        self.flush()
        self.conn.close()


# ── Scanner ───────────────────────────────────────────────────────────────

class PeanutScanner:
    def __init__(self, paths: list, recursive: bool = True, thumb_dir: str = "",
                 cache_dir: str = "", file_types: list = None):
        self.paths = [os.path.abspath(os.path.expanduser(p)) for p in paths]
        self.recursive = recursive
        self.thumb_dir = thumb_dir
        self.has_ffmpeg = ffmpeg_ok()
        self.cache = HashCache(
            os.path.join(cache_dir or os.path.expanduser("~/.peanut"), "cache.db")
        )
        # Build allowed extensions from selected categories
        if file_types:
            self.allowed_exts = set()
            for cat in file_types:
                if cat in FILE_CATEGORIES:
                    self.allowed_exts |= FILE_CATEGORIES[cat]
        else:
            self.allowed_exts = ALL_EXTS
        self.stats = {
            "total_files": 0, "images": 0, "videos": 0, "documents": 0,
            "audio": 0, "code": 0, "archives": 0, "other": 0,
            "exact_groups": 0, "similar_groups": 0,
            "total_groups": 0, "recoverable_bytes": 0,
            "cached_hits": 0,
        }

    def _emit(self, event: str, data: dict) -> dict:
        return {"event": event, "data": data}

    # ── Parallel hashing helpers ─────────────────────────────────────────
    # ThreadPoolExecutor works well here because the heavy lifting in
    # image_phash / video_phash / audio_fingerprint is in C extensions
    # (PIL, numpy, ffmpeg subprocess) that release the GIL. We submit
    # work to a worker pool, collect results in the main thread via
    # as_completed (so we can yield progress events from the generator),
    # and write to the cache from the main thread to keep SQLite writes
    # serialized.

    def _hash_one_image(self, f) -> Optional[str]:
        """Compute a phash for one image file (called from worker thread)."""
        try:
            return image_phash(f.path)
        except Exception:
            return None

    def _hash_one_video(self, f) -> Optional[str]:
        """Compute a video phash (called from worker thread)."""
        try:
            return video_phash(f.path)
        except Exception:
            return None

    def _hash_one_audio(self, f) -> Optional[str]:
        """Compute an audio fingerprint (called from worker thread)."""
        try:
            return audio_fingerprint(f.path)
        except Exception:
            return None

    def _parallel_compute(self, files: list, hash_fn, cache_data_fn,
                          phase_name: str, type_label: str) -> Generator:
        """
        Hash files in parallel using ThreadPoolExecutor, yielding progress
        events for live UI updates and writing results to the cache from
        the main thread.

        ``hash_fn(file_info) -> hash_str or None`` is called on workers.
        ``cache_data_fn(file_info) -> dict`` shapes the cache row.
        """
        if not files:
            return
        n = len(files)
        # Cap at 8 workers — beyond that, ffmpeg subprocesses thrash the disk
        # and SQLite contention starts to dominate.
        workers = max(1, min(8, os.cpu_count() or 4, n))
        yield self._emit("status", {
            "message": f"Hashing {n} {type_label} ({workers} workers)...",
            "phase": phase_name,
        })

        done = 0
        with ThreadPoolExecutor(max_workers=workers) as ex:
            future_to_file = {ex.submit(hash_fn, f): f for f in files}
            for fut in as_completed(future_to_file):
                f = future_to_file[fut]
                try:
                    h = fut.result()
                    if h:
                        f.phash = h
                        # Cache write from main thread (avoids SQLite locking issues)
                        self.cache.put(f.path, f.size, f.modified, cache_data_fn(f))
                except Exception:
                    pass
                done += 1
                if done % 20 == 0 or done == n:
                    yield self._emit("progress", {
                        "message": f"Hashed {done}/{n} {type_label}",
                        "current": done, "total": n,
                        "phase": phase_name,
                    })
        self.cache.flush()

    def scan_progressive(self) -> Generator:
        """Main scan pipeline. Yields events progressively."""
        yield self._emit("status", {"message": "Cataloging files...", "phase": "catalog"})

        # 1. Catalog
        file_paths = self._catalog()
        yield self._emit("status", {
            "message": f"Found {len(file_paths)} media files",
            "phase": "catalog",
            "total_files": len(file_paths),
        })

        if not file_paths:
            yield self._emit("complete", {"stats": self.stats})
            return

        # 2. Build file info with metadata
        yield self._emit("status", {"message": "Analyzing files...", "phase": "analyze"})
        files = []
        for i, fp in enumerate(file_paths):
            info = self._analyze_file(fp)
            if info:
                files.append(info)
            if (i + 1) % 50 == 0 or i == len(file_paths) - 1:
                yield self._emit("progress", {
                    "message": f"Analyzed {i+1}/{len(file_paths)} files",
                    "current": i + 1,
                    "total": len(file_paths),
                    "phase": "analyze",
                })

        self.stats["total_files"] = len(files)
        self.stats["images"] = sum(1 for f in files if f.file_type == "image")
        self.stats["videos"] = sum(1 for f in files if f.file_type == "video")
        self.stats["documents"] = sum(1 for f in files if f.file_type == "document")
        self.stats["audio"] = sum(1 for f in files if f.file_type == "audio")
        self.stats["code"] = sum(1 for f in files if f.file_type == "code")
        self.stats["archives"] = sum(1 for f in files if f.file_type == "archive")
        self.stats["other"] = sum(1 for f in files if f.file_type == "other")

        # 3. Exact duplicates
        yield self._emit("status", {"message": "Finding exact duplicates...", "phase": "exact"})
        group_id = 0
        exact_groups = []
        remaining_set = set()
        size_groups = defaultdict(list)
        for f in files:
            size_groups[f.size].append(f)
        candidates = {s: g for s, g in size_groups.items() if len(g) >= 2}
        total_candidates = sum(len(g) for g in candidates.values())
        hashed = 0
        matched = set()

        for size, group in candidates.items():
            md5_map = defaultdict(list)
            for f in group:
                if not f.md5:
                    f.md5 = md5_hash(f.path)
                hashed += 1
                if hashed % 100 == 0:
                    yield self._emit("progress", {
                        "message": f"Hashing {hashed}/{total_candidates} files...",
                        "current": hashed, "total": total_candidates, "phase": "exact",
                    })
                if f.md5:
                    md5_map[f.md5].append(f)
            for md5, dupes in md5_map.items():
                if len(dupes) >= 2:
                    for d in dupes:
                        d.similarity = 100.0
                    rec = pick_best(dupes)
                    g = DuplicateGroup(group_id=group_id, match_type="exact",
                                       files=dupes, recommended_keep=rec)
                    group_id += 1
                    if self.thumb_dir:
                        for f in g.files:
                            f.thumbnail = make_thumbnail(f.path, self.thumb_dir, f.file_type)
                    exact_groups.append(g)
                    yield self._emit("group", g.to_dict())
                    for d in dupes:
                        matched.add(d.path)
                    self.stats["recoverable_bytes"] += sum(
                        d.size for d in dupes if d.path != rec)

        remaining = [f for f in files if f.path not in matched]
        self.stats["exact_groups"] = len(exact_groups)

        # 4. Similar images
        images = [f for f in remaining if f.file_type == "image"]
        if images:
            # Parallel phash computation (skips files already cached)
            images_to_hash = [f for f in images if not f.phash]
            yield from self._parallel_compute(
                images_to_hash, self._hash_one_image,
                lambda f: {
                    "md5": f.md5, "phash": f.phash, "sharpness": f.sharpness,
                    "exif": f.exif, "source": f.source, "compression": f.compression,
                    "width": f.width, "height": f.height,
                },
                "hash_images", "images",
            )
            yield self._emit("status", {
                "message": f"Comparing {len(images)} images...",
                "phase": "similar_images",
            })
            sim_img_groups = self._find_similar_images(images)
            for g in sim_img_groups:
                g.group_id = group_id
                group_id += 1
                if self.thumb_dir:
                    for f in g.files:
                        f.thumbnail = make_thumbnail(f.path, self.thumb_dir, f.file_type)
                yield self._emit("group", g.to_dict())
            self.stats["similar_groups"] += len(sim_img_groups)
            matched_img_paths = {f.path for g in sim_img_groups for f in g.files}
            remaining = [f for f in remaining if f.path not in matched_img_paths]

        # 5. Similar videos
        videos = [f for f in remaining if f.file_type == "video"]
        if videos and self.has_ffmpeg:
            videos_to_hash = [f for f in videos if not f.phash]
            yield from self._parallel_compute(
                videos_to_hash, self._hash_one_video,
                lambda f: {
                    "md5": f.md5, "phash": f.phash,
                    "width": f.width, "height": f.height,
                },
                "hash_videos", "videos",
            )
            yield self._emit("status", {
                "message": f"Comparing {len(videos)} videos...",
                "phase": "similar_videos",
            })
            sim_vid_groups = self._find_similar_videos(videos)
            for g in sim_vid_groups:
                g.group_id = group_id
                group_id += 1
                if self.thumb_dir:
                    for f in g.files:
                        f.thumbnail = make_thumbnail(f.path, self.thumb_dir, f.file_type)
                yield self._emit("group", g.to_dict())
            self.stats["similar_groups"] += len(sim_vid_groups)

        # 6. Similar documents
        docs = [f for f in remaining if f.file_type == "document"]
        if docs:
            yield self._emit("status", {
                "message": f"Comparing {len(docs)} documents...",
                "phase": "similar_docs",
            })
            doc_groups = self._find_similar_documents(docs)
            for g in doc_groups:
                g.group_id = group_id
                group_id += 1
                yield self._emit("group", g.to_dict())
            self.stats["similar_groups"] += len(doc_groups)
            matched_doc_paths = {f.path for g in doc_groups for f in g.files}
            remaining = [f for f in remaining if f.path not in matched_doc_paths]

        # 7. Similar code files (same SimHash approach as documents)
        code_files = [f for f in remaining if f.file_type == "code"]
        if code_files:
            yield self._emit("status", {
                "message": f"Comparing {len(code_files)} code files...",
                "phase": "similar_code",
            })
            code_groups = self._find_similar_documents(code_files)
            for g in code_groups:
                g.group_id = group_id
                group_id += 1
                yield self._emit("group", g.to_dict())
            self.stats["similar_groups"] += len(code_groups)
            matched_code_paths = {f.path for g in code_groups for f in g.files}
            remaining = [f for f in remaining if f.path not in matched_code_paths]

        # 8. Similar audio files
        audio_files = [f for f in remaining if f.file_type == "audio"]
        if audio_files and self.has_ffmpeg:
            audio_to_hash = [f for f in audio_files if not f.phash]
            yield from self._parallel_compute(
                audio_to_hash, self._hash_one_audio,
                lambda f: {
                    "md5": f.md5, "phash": f.phash,
                    "compression": f.compression,
                },
                "hash_audio", "audio files",
            )
            yield self._emit("status", {
                "message": f"Comparing {len(audio_files)} audio files...",
                "phase": "similar_audio",
            })
            audio_groups = self._find_similar_audio(audio_files)
            for g in audio_groups:
                g.group_id = group_id
                group_id += 1
                yield self._emit("group", g.to_dict())
            self.stats["similar_groups"] += len(audio_groups)

        # Done
        self.stats["total_groups"] = group_id
        yield self._emit("complete", {"stats": self.stats})
        self.cache.close()

    def _catalog(self) -> list:
        paths = []
        seen = set()
        for scan_path in self.paths:
            if not os.path.exists(scan_path):
                continue
            if os.path.isfile(scan_path):
                ext = Path(scan_path).suffix.lower()
                if ext in self.allowed_exts:
                    ap = os.path.abspath(scan_path)
                    if ap not in seen:
                        seen.add(ap)
                        paths.append(ap)
                continue
            walker = os.walk(scan_path) if self.recursive else \
                [(scan_path, [], os.listdir(scan_path))]
            for root, dirs, fnames in walker:
                dirs[:] = [d for d in dirs if not d.startswith(".")]
                for fn in fnames:
                    if fn.startswith("."):
                        continue
                    ext = Path(fn).suffix.lower()
                    if ext in self.allowed_exts:
                        fp = os.path.abspath(os.path.join(root, fn))
                        if fp not in seen:
                            seen.add(fp)
                            paths.append(fp)
        return paths

    def _analyze_file(self, path: str) -> Optional[FileInfo]:
        try:
            stat = os.stat(path)
            if stat.st_size < MIN_FILE_SIZE:
                return None
            ext = Path(path).suffix.lower()
            ft = ("image" if ext in IMAGE_EXTS else
                  "video" if ext in VIDEO_EXTS else
                  "audio" if ext in AUDIO_EXTS else
                  "code" if ext in CODE_EXTS else
                  "archive" if ext in ARCHIVE_EXTS else
                  "document" if ext in DOC_EXTS else
                  "other")

            info = FileInfo(path=path, size=stat.st_size, file_type=ft,
                            modified=stat.st_mtime)

            # Check cache
            cached = self.cache.get(path, stat.st_size, stat.st_mtime)
            if cached:
                self.stats["cached_hits"] += 1
                info.md5 = cached["md5"]
                info.phash = cached["phash"]
                info.sharpness = cached["sharpness"]
                info.exif = cached["exif"]
                info.source = cached["source"]
                info.compression = cached["compression"]
                info.width = cached["width"]
                info.height = cached["height"]
                info.text_hash = cached.get("text_hash", "")
                info.word_count = cached.get("word_count", 0)
                info.text_preview = cached.get("text_preview", "")
            else:
                if ft == "image":
                    w, h = image_dimensions(path)
                    info.width, info.height = w, h
                    info.sharpness = compute_sharpness(path)
                    info.exif = extract_exif(path)
                    info.source = detect_source(path, info.exif)
                    info.compression = detect_compression(path)
                elif ft == "video" and self.has_ffmpeg:
                    info.duration = video_duration(path)
                    w, h = video_dimensions(path)
                    info.width, info.height = w, h
                    info.source = detect_source(path, {})
                    info.compression = "lossy"
                elif ft == "audio" and self.has_ffmpeg:
                    info.duration = audio_duration(path)
                    info.source = detect_source(path, {})
                    info.compression = detect_compression(path)
                elif ft == "document":
                    text = extract_text_from_file(path)
                    info.text_hash = compute_simhash(text)
                    info.word_count = len(text.split()) if text else 0
                    info.text_preview = text[:300].strip() if text else ""
                    info.source = detect_source(path, {})
                elif ft == "code":
                    text = extract_text_from_file(path)
                    info.text_hash = compute_simhash(text)
                    info.word_count = len(text.split()) if text else 0
                    info.text_preview = text[:300].strip() if text else ""
                elif ft == "archive":
                    info.md5 = archive_content_hash(path)
                # ft == "other": no extra analysis — exact-dedupe via md5 only
                # (md5 is computed during the exact-match phase, not here).

                # Cache it
                self.cache.put(path, stat.st_size, stat.st_mtime, {
                    "md5": info.md5, "phash": info.phash,
                    "sharpness": info.sharpness, "exif": info.exif,
                    "source": info.source, "compression": info.compression,
                    "width": info.width, "height": info.height,
                    "text_hash": info.text_hash, "word_count": info.word_count,
                    "text_preview": info.text_preview,
                })

            info.dimensions = f"{info.width}x{info.height}" if info.width else ""
            if info.width and info.height and info.size:
                info.bits_per_pixel = (info.size * 8) / (info.width * info.height)
            info.quality_score = compute_quality_score(info)
            return info
        except (OSError, PermissionError):
            return None

    def _find_exact(self, files: list) -> tuple:
        size_groups = defaultdict(list)
        for f in files:
            size_groups[f.size].append(f)

        candidates = {s: g for s, g in size_groups.items() if len(g) >= 2}
        groups = []
        matched = set()

        for size, group in candidates.items():
            md5_map = defaultdict(list)
            for f in group:
                if not f.md5:
                    f.md5 = md5_hash(f.path)
                if f.md5:
                    md5_map[f.md5].append(f)
            for md5, dupes in md5_map.items():
                if len(dupes) >= 2:
                    for d in dupes:
                        d.similarity = 100.0
                    rec = pick_best(dupes)
                    groups.append(DuplicateGroup(
                        group_id=0, match_type="exact",
                        files=dupes, recommended_keep=rec,
                    ))
                    for d in dupes:
                        matched.add(d.path)
                    self.stats["recoverable_bytes"] += sum(
                        d.size for d in dupes if d.path != rec)

        remaining = [f for f in files if f.path not in matched]
        return groups, remaining

    def _find_similar_images(self, images: list) -> list:
        # Compute phashes (skip if cached). NOTE: parallel hashing is done
        # in scan_progressive before this is called, so this loop is a
        # safety net for files whose phash wasn't pre-computed.
        for f in images:
            if not f.phash:
                h = image_phash(f.path)
                if h:
                    f.phash = h
                    self.cache.put(f.path, f.size, f.modified, {
                        "md5": f.md5, "phash": f.phash, "sharpness": f.sharpness,
                        "exif": f.exif, "source": f.source, "compression": f.compression,
                        "width": f.width, "height": f.height,
                    })
        self.cache.flush()

        hashed = [f for f in images if f.phash]
        if len(hashed) < 2:
            return []

        # Vectorized Hamming-distance computation.
        # Standard imagehash phash is 16 hex chars (64 bits).
        packed = _hex_to_bits_array([f.phash for f in hashed], expected_hex_len=16)
        pair_list = _vectorized_pairs(packed, IMAGE_THRESHOLD, divisor=1)

        # Union-find clustering
        parent = {f.path: f.path for f in hashed}
        distances = {}  # (path1, path2) -> distance

        def find(x):
            while parent[x] != x:
                parent[x] = parent[parent[x]]
                x = parent[x]
            return x

        def union(a, b):
            ra, rb = find(a), find(b)
            if ra != rb:
                parent[ra] = rb

        for i, j, dist in pair_list:
            pi, pj = hashed[i].path, hashed[j].path
            union(pi, pj)
            distances[(pi, pj)] = dist
            distances[(pj, pi)] = dist

        clusters = defaultdict(list)
        for f in hashed:
            clusters[find(f.path)].append(f)

        groups = []
        for root, cluster in clusters.items():
            if len(cluster) >= 2:
                # Compute similarity for each file relative to the best
                rec = pick_best(cluster)
                best_file = next(f for f in cluster if f.path == rec)
                for f in cluster:
                    if f.path == rec:
                        f.similarity = 100.0
                    else:
                        d = distances.get((f.path, rec))
                        if d is None:
                            d = hamming(f.phash, best_file.phash)
                        f.similarity = similarity_pct(d)
                groups.append(DuplicateGroup(
                    group_id=0, match_type="similar",
                    files=cluster, recommended_keep=rec,
                ))
                self.stats["recoverable_bytes"] += sum(
                    f.size for f in cluster if f.path != rec)
        return groups

    def _find_similar_videos(self, videos: list) -> list:
        for f in videos:
            if not f.phash:
                h = video_phash(f.path)
                if h:
                    f.phash = h
                    self.cache.put(f.path, f.size, f.modified, {
                        "md5": f.md5, "phash": f.phash,
                        "width": f.width, "height": f.height,
                    })
        self.cache.flush()

        hashed = [f for f in videos if f.phash]
        if len(hashed) < 2:
            return []

        # Vectorized multi-segment Hamming.
        # video_phash emits "h1;h2;...;hN" — N keyframes, each 16x16 = 256 bits = 64 hex chars.
        # We bucket by segment count: same-count videos can be vectorized together.
        # Cross-bucket pairs (rare — happens only if ffmpeg missed keyframes) use Python fallback.
        pair_list = self._multisegment_pairs(
            hashed, threshold=VIDEO_THRESHOLD,
            seg_hex_len=64,            # 16x16 phash = 256 bits = 64 hex chars
            python_fallback=video_hamming,
        )

        # Apply duration prefilter (matches original behavior)
        def duration_ok(i, j):
            d1, d2 = hashed[i].duration, hashed[j].duration
            if d1 > 0 and d2 > 0 and min(d1, d2) / max(d1, d2) < 0.9:
                return False
            return True
        pair_list = [(i, j, d) for i, j, d in pair_list if duration_ok(i, j)]

        # Union-find clustering
        parent = {f.path: f.path for f in hashed}
        distances = {}

        def find(x):
            while parent[x] != x:
                parent[x] = parent[parent[x]]
                x = parent[x]
            return x

        def union(a, b):
            ra, rb = find(a), find(b)
            if ra != rb:
                parent[ra] = rb

        for i, j, dist in pair_list:
            pi, pj = hashed[i].path, hashed[j].path
            union(pi, pj)
            distances[(pi, pj)] = dist
            distances[(pj, pi)] = dist

        clusters = defaultdict(list)
        for f in hashed:
            clusters[find(f.path)].append(f)

        groups = []
        for root, cluster in clusters.items():
            if len(cluster) >= 2:
                rec = pick_best(cluster)
                best = next(f for f in cluster if f.path == rec)
                for f in cluster:
                    if f.path == rec:
                        f.similarity = 100.0
                    else:
                        d = distances.get((f.path, rec))
                        if d is None:
                            d = video_hamming(f.phash, best.phash)
                        f.similarity = similarity_pct(d)
                groups.append(DuplicateGroup(
                    group_id=0, match_type="similar",
                    files=cluster, recommended_keep=rec,
                ))
                self.stats["recoverable_bytes"] += sum(
                    f.size for f in cluster if f.path != rec)
        return groups

    def _multisegment_pairs(self, hashed: list, threshold: float,
                            seg_hex_len: int, python_fallback) -> list:
        """
        Find below-threshold pairs for multi-segment hashes (videos, audio).
        Hashes are ``;``-separated segments. We bucket by segment count and
        vectorize within each bucket; cross-bucket pairs (rare) fall back to
        the supplied Python distance function.
        """
        # Bucket indices by segment count
        groups_by_segs = defaultdict(list)
        for idx, f in enumerate(hashed):
            nsegs = len(f.phash.split(";"))
            groups_by_segs[nsegs].append(idx)

        pairs = []
        # Vectorize within each bucket
        for nsegs, indices in groups_by_segs.items():
            if len(indices) < 2 or nsegs == 0:
                continue
            # Concatenate all segments per video into one big hex string
            concat = []
            for i in indices:
                segs = hashed[i].phash.split(";")
                # Normalize each segment length, then join
                fixed = "".join(s.zfill(seg_hex_len)[:seg_hex_len] for s in segs)
                concat.append(fixed)
            packed = _hex_to_bits_array(concat, expected_hex_len=nsegs * seg_hex_len)
            for i_local, j_local, dist in _vectorized_pairs(packed, threshold, divisor=nsegs):
                pairs.append((indices[i_local], indices[j_local], dist))

        # Cross-bucket comparisons (videos with different segment counts) — Python fallback.
        seg_counts = sorted(groups_by_segs.keys())
        for k1 in range(len(seg_counts)):
            for k2 in range(k1 + 1, len(seg_counts)):
                for i in groups_by_segs[seg_counts[k1]]:
                    for j in groups_by_segs[seg_counts[k2]]:
                        dist = python_fallback(hashed[i].phash, hashed[j].phash)
                        if dist <= threshold:
                            lo, hi = (i, j) if i < j else (j, i)
                            pairs.append((lo, hi, dist))
        return pairs

    def _find_similar_documents(self, docs: list) -> list:
        """Find similar documents via SimHash text fingerprinting."""
        hashed = [f for f in docs if f.text_hash]
        if len(hashed) < 2:
            return []

        # Vectorized Hamming-distance computation.
        # SimHash is 128 bits = 32 hex chars (compute_simhash uses hash_bits=128 by default).
        packed = _hex_to_bits_array([f.text_hash for f in hashed], expected_hex_len=32)
        pair_list = _vectorized_pairs(packed, DOC_THRESHOLD, divisor=1)

        parent = {f.path: f.path for f in hashed}
        distances = {}

        def find(x):
            while parent[x] != x:
                parent[x] = parent[parent[x]]
                x = parent[x]
            return x

        def union(a, b):
            ra, rb = find(a), find(b)
            if ra != rb:
                parent[ra] = rb

        for i, j, dist in pair_list:
            pi, pj = hashed[i].path, hashed[j].path
            union(pi, pj)
            distances[(pi, pj)] = dist
            distances[(pj, pi)] = dist

        clusters = defaultdict(list)
        for f in hashed:
            clusters[find(f.path)].append(f)

        groups = []
        for root, cluster in clusters.items():
            if len(cluster) >= 2:
                # For documents, prefer the largest/newest file
                rec = max(cluster, key=lambda f: (f.size, f.modified)).path
                for f in cluster:
                    if f.path == rec:
                        f.similarity = 100.0
                    else:
                        d = distances.get((f.path, rec))
                        if d is None:
                            best = next(x for x in cluster if x.path == rec)
                            d = simhash_distance(f.text_hash, best.text_hash)
                        f.similarity = similarity_pct(d)
                groups.append(DuplicateGroup(
                    group_id=0, match_type="similar",
                    files=cluster, recommended_keep=rec,
                ))
                self.stats["recoverable_bytes"] += sum(
                    f.size for f in cluster if f.path != rec)
        return groups

    def _find_similar_audio(self, audio_files: list) -> list:
        """Find similar audio files via audio fingerprinting."""
        # Compute fingerprints
        for f in audio_files:
            if not f.phash:
                h = audio_fingerprint(f.path)
                if h:
                    f.phash = h
                    self.cache.put(f.path, f.size, f.modified, {
                        "md5": f.md5, "phash": f.phash,
                        "compression": f.compression,
                    })
        self.cache.flush()

        hashed = [f for f in audio_files if f.phash]
        if len(hashed) < 2:
            return []

        # Vectorized multi-segment Hamming for audio.
        # audio_fingerprint emits "h1;h2;...;h10" — 10 segments, each 32 bits = 8 hex chars.
        pair_list = self._multisegment_pairs(
            hashed, threshold=AUDIO_THRESHOLD,
            seg_hex_len=8,             # 32 bits = 8 hex chars per segment
            python_fallback=audio_hash_distance,
        )

        # Apply duration prefilter (matches original behavior)
        def duration_ok(i, j):
            d1, d2 = hashed[i].duration, hashed[j].duration
            if d1 > 0 and d2 > 0 and min(d1, d2) / max(d1, d2) < 0.9:
                return False
            return True
        pair_list = [(i, j, d) for i, j, d in pair_list if duration_ok(i, j)]

        parent = {f.path: f.path for f in hashed}
        distances = {}

        def find(x):
            while parent[x] != x:
                parent[x] = parent[parent[x]]
                x = parent[x]
            return x

        def union(a, b):
            ra, rb = find(a), find(b)
            if ra != rb:
                parent[ra] = rb

        for i, j, dist in pair_list:
            pi, pj = hashed[i].path, hashed[j].path
            union(pi, pj)
            distances[(pi, pj)] = dist
            distances[(pj, pi)] = dist

        clusters = defaultdict(list)
        for f in hashed:
            clusters[find(f.path)].append(f)

        groups = []
        for root, cluster in clusters.items():
            if len(cluster) >= 2:
                # Prefer lossless, then largest
                rec = max(cluster, key=lambda f: (
                    1 if f.compression == "lossless" else 0, f.size, f.modified
                )).path
                for f in cluster:
                    if f.path == rec:
                        f.similarity = 100.0
                    else:
                        d = distances.get((f.path, rec))
                        if d is None:
                            best = next(x for x in cluster if x.path == rec)
                            d = audio_hash_distance(f.phash, best.phash)
                        f.similarity = similarity_pct(d)
                groups.append(DuplicateGroup(
                    group_id=0, match_type="similar",
                    files=cluster, recommended_keep=rec,
                ))
                self.stats["recoverable_bytes"] += sum(
                    f.size for f in cluster if f.path != rec)
        return groups


# ──────────────────────────────────────────────────────────────────────────
#                    For whoever made it to the bottom.
# ──────────────────────────────────────────────────────────────────────────
#
# The thing that's going to change your life is not a habit, a morning
# routine, or a better system. It's one decision made from the part of
# you that is done waiting for conditions to be perfect. One ugly,
# half-ready, terrifying move made before you feel ready to make it.
#
# You think I was ready to ship my app? I wasn't. I spent months building
# that thing, and every single night before bed I'd tell myself "tomorrow's
# the day." I could not wait to lock in and finally ship it. And then
# tomorrow would come, and I would find a reason not to.
#
# Everything you admire about someone else's life started with them just
# doing the thing before they were ready. And that moment looked nothing
# like a beginning. It looked like a bad idea. It looked like the wrong
# time. It looked like something they'd probably regret.
#
# They did it anyway. That's the whole story.
#
#                                                          — cody browne
# ──────────────────────────────────────────────────────────────────────────
